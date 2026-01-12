from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

theme = 'corporate' # Default to corporate as no theme was provided

if theme == 'corporate':
    bg_color = RGBColor(0xF2, 0xF2, 0xF2)
    title_color = RGBColor(0x1F, 0x49, 0x7D)
    bullet_color_level0 = RGBColor(0x00, 0x00, 0x00)
    bullet_color_level1 = RGBColor(0x40, 0x40, 0x40)
elif theme == 'minimal':
    bg_color = RGBColor(0xFF, 0xFF, 0xFF)
    title_color = RGBColor(0x33, 0x33, 0x33)
    bullet_color_level0 = RGBColor(0x00, 0x00, 0x00)
    bullet_color_level1 = RGBColor(0x66, 0x66, 0x66)
elif theme == 'dark':
    bg_color = RGBColor(0x22, 0x22, 0x22)
    title_color = RGBColor(0xFF, 0xFF, 0xFF)
    bullet_color_level0 = RGBColor(0xE0, 0xE0, 0xE0)
    bullet_color_level1 = RGBColor(0xBB, 0xBB, 0xBB)
elif theme == 'academic':
    bg_color = RGBColor(0xEE, 0xEE, 0xEE)
    title_color = RGBColor(0x00, 0x00, 0x80) # Navy blue
    bullet_color_level0 = RGBColor(0x00, 0x00, 0x00)
    bullet_color_level1 = RGBColor(0x33, 0x33, 0x33)
else: # Default to corporate if invalid theme
    bg_color = RGBColor(0xF2, 0xF2, 0xF2)
    title_color = RGBColor(0x1F, 0x49, 0x7D)
    bullet_color_level0 = RGBColor(0x00, 0x00, 0x00)
    bullet_color_level1 = RGBColor(0x40, 0x40, 0x40)

def set_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, topic, subtitle, title_color, bg_color):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, bg_color)

    title = slide.shapes.title
    title.text = topic
    title.text_frame.paragraphs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.name = 'Arial'
    subtitle_shape.paragraphs[0].font.size = Pt(28)
    subtitle_shape.paragraphs[0].font.color.rgb = bullet_color_level0
    subtitle_shape.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

def add_agenda_slide(prs, slide_titles, title_color, bullet_color_level0, bg_color):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, bg_color)

    title = slide.shapes.title
    title.text = "Agenda"
    title.text_frame.paragraphs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    for title_text in slide_titles:
        p = tf.add_paragraph()
        p.text = title_text
        p.level = 0
        p.font.name = 'Arial'
        p.font.size = Pt(24)
        p.font.color.rgb = bullet_color_level0

def add_content_slide(prs, slide_title, content_list, title_color, bullet_color_level0, bullet_color_level1, bg_color):
    MAX_BULLETS = 7
    all_content_chunks = [content_list[i:i + MAX_BULLETS] for i in range(0, len(content_list), MAX_BULLETS)]

    for i, chunk in enumerate(all_content_chunks):
        current_slide_title = slide_title
        if i > 0:
            current_slide_title += " (contd.)"

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, bg_color)

        title = slide.shapes.title
        title.text = current_slide_title
        title.text_frame.paragraphs[0].font.name = 'Arial'
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = title_color
        title.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        for item in chunk:
            level = 0
            text_to_add = item
            if isinstance(item, tuple): # Handle nested bullets
                level = item[0]
                text_to_add = item[1]

            if len(text_to_add) > 110:
                text_to_add = text_to_add[:107] + "..."

            p = tf.add_paragraph()
            p.text = text_to_add
            p.level = level
            p.font.name = 'Arial'
            if level == 0:
                p.font.size = Pt(24)
                p.font.color.rgb = bullet_color_level0
            else:
                p.font.size = Pt(20)
                p.font.color.rgb = bullet_color_level1

def add_summary_slide(prs, topic, summary_points, title_color, bullet_color_level0, bg_color):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, bg_color)

    title = slide.shapes.title
    title.text = f"Summary: {topic}"
    title.text_frame.paragraphs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    for point in summary_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.name = 'Arial'
        p.font.size = Pt(24)
        p.font.color.rgb = bullet_color_level0

def add_qa_slide(prs, title_color, bg_color):
    slide_layout = prs.slide_layouts[5] # Blank layout for centered text
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, bg_color)

    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1.5)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Any Questions?"
    
    p = tf.paragraphs[0]
    p.alignment = 2 # Center alignment
    p.font.name = 'Arial'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = title_color

def add_thank_you_slide(prs, title_color, bg_color):
    slide_layout = prs.slide_layouts[5] # Blank layout for centered text
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, bg_color)

    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1.5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Thank You!"
    
    p = tf.paragraphs[0]
    p.alignment = 2 # Center alignment
    p.font.name = 'Arial'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = title_color

topic = "Introduction to Artificial Intelligence"
subtitle = "An Informative Presentation"

slide_data = [
    {
        "title": "What is Artificial Intelligence?",
        "content": [
            "Definition: Simulation of human intelligence in machines.",
            "Key Goal: Enable machines to perform tasks typically requiring human intellect.",
            "Sub-fields: Machine Learning, Deep Learning, Natural Language Processing.",
            "Evolution: From rule-based systems to complex neural networks."
        ]
    },
    {
        "title": "Core Concepts of AI",
        "content": [
            "Machine Learning (ML): Algorithms learning from data without explicit programming.",
            "Deep Learning (DL): Subset of ML using neural networks with many layers.",
            "Natural Language Processing (NLP): Understanding and generating human language.",
            "Computer Vision (CV): Enabling machines to \"see\" and interpret visual data.",
            "Robotics: Integrating AI into physical machines for automation."
        ]
    },
    {
        "title": "Types of AI",
        "content": [
            "Narrow AI (ANI): Task-specific intelligence (e.g., Siri, self-driving cars).",
            "General AI (AGI): Human-level intelligence across multiple tasks (hypothetical).",
            "Superintelligence (ASI): AI surpassing human intellect (speculative future).",
            "Reactive Machines: Basic AI, no memory, only react (e.g., chess AI).",
            "Limited Memory AI: Stores past data for a short time (e.g., autonomous vehicles)."
        ]
    },
    {
        "title": "Applications of AI",
        "content": [
            "Healthcare: Disease diagnosis, drug discovery, personalized medicine.",
            "Finance: Fraud detection, algorithmic trading, credit scoring.",
            "Retail: Recommendation systems, inventory management, customer service bots.",
            "Transportation: Self-driving cars, traffic optimization, logistics.",
            "Education: Personalized learning, intelligent tutoring systems.",
            "Entertainment: Content recommendations, game AI, special effects."
        ]
    },
    {
        "title": "How AI Works (Simplified)",
        "content": [
            "Data Collection: Gathering vast amounts of relevant data.",
            "Algorithm Selection: Choosing the appropriate AI model (e.g., neural network, decision tree).",
            "Training: Feeding data to the model to learn patterns and make predictions.",
            "Evaluation: Testing the model's accuracy and performance.",
            "Deployment: Integrating the trained model into real-world applications.",
            "Continuous Learning: Models can improve over time with new data."
        ]
    },
    {
        "title": "Challenges and Ethical Considerations",
        "content": [
            "Bias in Data: AI models can perpetuate and amplify existing biases.",
            "Job Displacement: Automation raises concerns about future employment.",
            "Privacy Concerns: Extensive data collection for AI models.",
            "Accountability: Who is responsible when AI makes mistakes?",
            "Explainability (XAI): Difficulty in understanding how complex models make decisions.",
            "Security Risks: AI systems can be vulnerable to attacks."
        ]
    },
    {
        "title": "The Future of AI",
        "content": [
            "Continued Integration: AI will become more pervasive in daily life.",
            "Advancements in AGI: Research continues towards more general intelligence.",
            "Human-AI Collaboration: Focus on augmenting human capabilities.",
            "Personalized Experiences: More tailored services and products.",
            "New Discoveries: AI accelerating scientific research and innovation.",
            "Regulatory Frameworks: Increasing need for ethical guidelines and laws."
        ]
    }
]

agenda_titles = [s["title"] for s in slide_data]

add_title_slide(prs, topic, subtitle, title_color, bg_color)
add_agenda_slide(prs, agenda_titles, title_color, bullet_color_level0, bg_color)

for data in slide_data:
    add_content_slide(prs, data["title"], data["content"], title_color, bullet_color_level0, bullet_color_level1, bg_color)

summary_points = [
    "AI is transforming industries by simulating human intelligence.",
    "Key areas include ML, DL, NLP, and Computer Vision.",
    "Applications span healthcare, finance, retail, and more.",
    "Ethical considerations like bias and privacy are crucial for responsible development.",
    "The future involves deeper integration and collaborative human-AI systems."
]
add_summary_slide(prs, topic, summary_points, title_color, bullet_color_level0, bg_color)

add_qa_slide(prs, title_color, bg_color)
add_thank_you_slide(prs, title_color, bg_color)

prs.save("data.pptx")